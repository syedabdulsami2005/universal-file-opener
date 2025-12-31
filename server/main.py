import magic
import pandas as pd
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
import nbformat
from nbconvert import HTMLExporter
import mammoth
import io
import ebooklib
from ebooklib import epub
from pptx import Presentation
from bs4 import BeautifulSoup

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/detect-and-convert")
async def detect_and_convert(file: UploadFile = File(...)):
    content = await file.read()
    filename = file.filename.lower()
    
    # 1. Identify Mime Type
    try:
        mime = magic.from_buffer(content, mime=True)
    except:
        mime = "application/octet-stream"

    buffer = io.BytesIO(content)

    # --- A. MICROSOFT OFFICE ---
    
    # Word (.docx)
    if filename.endswith('.docx'):
        try:
            result = mammoth.convert_to_html(buffer)
            return {"type": "html_doc", "content": f'<div style="padding:20px; bg:white;">{result.value}</div>', "mime": mime}
        except Exception as e: return {"error": str(e)}

    # PowerPoint (.pptx) - Extract text from slides
    if filename.endswith('.pptx'):
        try:
            prs = Presentation(buffer)
            html_content = '<div style="padding:20px; background:white;">'
            for i, slide in enumerate(prs.slides):
                html_content += f'<div style="margin-bottom:20px; border:1px solid #ccc; padding:10px;"><h3>Slide {i+1}</h3>'
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        html_content += f'<p>{shape.text}</p>'
                html_content += '</div>'
            html_content += '</div>'
            return {"type": "html_doc", "content": html_content, "mime": mime}
        except Exception as e: return {"error": str(e)}

    # Excel / CSV (.xlsx, .csv)
    if filename.endswith(('.xlsx', '.xls', '.csv', '.parquet')):
        try:
            if filename.endswith('.csv'): df = pd.read_csv(buffer)
            elif filename.endswith('.parquet'): df = pd.read_parquet(buffer)
            else: df = pd.read_excel(buffer)
            return {"type": "html_table", "content": df.to_html(classes="min-w-full bg-white border"), "mime": mime}
        except Exception as e: return {"error": str(e)}

    # --- B. EBOOKS ---
    if filename.endswith('.epub'):
        try:
            # Basic EPUB text extraction (simplified for cloud)
            # Note: Full EPUB rendering is complex, this extracts text chapters
            book = epub.read_epub(io.BytesIO(content))
            html_out = '<div style="padding:20px; background:white;">'
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    html_out += str(soup.body) if soup.body else str(soup)
            html_out += '</div>'
            return {"type": "html_doc", "content": html_out, "mime": mime}
        except Exception as e: return {"error": str(e)}

    # --- C. NOTEBOOKS ---
    if filename.endswith('.ipynb'):
        try:
            nb = nbformat.reads(content.decode('utf-8'), as_version=4)
            html_exporter = HTMLExporter()
            (body, _) = html_exporter.from_notebook_node(nb)
            return {"type": "html_doc", "content": body, "mime": mime}
        except Exception as e: return {"error": str(e)}

    # --- D. FALLBACK ---
    # Crucial: Return specific type so frontend knows NOT to try and render text
    return {"type": "unknown", "mime": mime}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)